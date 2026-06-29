#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Genera la presentación ejecutiva (PPTX) para Conexión Talento.

v2 — reestructurada tras el Consejo LLM #01:
  · arco "empecemos por una victoria" (Fase 0 = la decisión del viernes)
  · roadmap completo reposicionado como horizonte/anexo
  · "SÍ/retamos" → "Sí, y"; nueva slide de ROI en dólares de Virginia; demo viva
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- Paleta ----------
NAVY  = RGBColor(0x10,0x2A,0x43)
NAVY2 = RGBColor(0x0C,0x22,0x36)
TEAL  = RGBColor(0x1C,0x7C,0x8C)
TEALL = RGBColor(0x2E,0x9C,0xAD)
DEEP  = RGBColor(0x14,0x50,0x6B)
GOLD  = RGBColor(0xC9,0xA2,0x27)
INK   = RGBColor(0x1A,0x27,0x33)
MUTED = RGBColor(0x5B,0x6B,0x7B)
LINE  = RGBColor(0xD9,0xE2,0xEC)
SOFT  = RGBColor(0xF5,0xF7,0xFA)
WHITE = RGBColor(0xFF,0xFF,0xFF)
FONT  = "Arial"
FIRM  = "‹TU CONSULTORA›"   # ← reemplazar por el nombre real de la firma

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
ML = Inches(0.7)                      # margen izq
CW = Inches(13.333 - 1.4)            # ancho de contenido

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, l, t, w, h, fill=None, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, l, t, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = line_w or Pt(1)
    return sp

def tb(s, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    b = s.shapes.add_textbox(l, t, w, h); tf = b.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    return tf

def para(tf, runs, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT,
         space_before=0, space_after=4, bullet=False, leading=1.08, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before); p.space_after = Pt(space_after)
    try: p.line_spacing = leading
    except Exception: pass
    if isinstance(runs, str):
        runs = [(runs, bold, color, size)]
    if bullet:
        r = p.add_run(); r.text = "•  "
        r.font.size = Pt(size); r.font.color.rgb = TEAL; r.font.bold = True; r.font.name = FONT
    for item in runs:
        text, b, c, sz = (item + (bold, color, size))[:4] if isinstance(item, tuple) else (item, bold, color, size)
        r = p.add_run(); r.text = text
        r.font.size = Pt(sz); r.font.bold = b; r.font.color.rgb = c; r.font.name = FONT
    return p

def footer(s, n, light=False):
    c = RGBColor(0x9F,0xB4,0xC6) if light else RGBColor(0xA8,0xB4,0xC0)
    tf = tb(s, ML, Inches(7.04), Inches(8), Inches(0.32))
    para(tf, "Confidencial · Conexión Talento", size=8, color=c, first=True)
    tf2 = tb(s, Inches(13.333-4.7), Inches(7.04), Inches(4), Inches(0.32))
    para(tf2, f"{FIRM}   ·   {n}", size=8, color=c, align=PP_ALIGN.RIGHT, first=True)

def title_bar(s, kicker, title, n):
    tf = tb(s, ML, Inches(0.42), CW, Inches(0.3))
    para(tf, kicker.upper(), size=10.5, color=TEAL, bold=True, first=True)
    tf2 = tb(s, ML, Inches(0.72), CW, Inches(0.7))
    para(tf2, title, size=25, color=NAVY, bold=True, first=True, leading=1.0)
    rect(s, ML, Inches(1.46), Inches(1.6), Pt(3), fill=TEAL)
    footer(s, n)
    return Inches(1.75)   # y de inicio de contenido

def style_table(tbl, header_fill=NAVY, header_color=WHITE, body_size=11.5,
                header_size=11.5, zebra=True, align_first_left=True):
    for ci, cell in enumerate(tbl.rows[0].cells):
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
        cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
        cell.margin_left = Pt(6); cell.margin_right = Pt(6)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(header_size); r.font.bold = True
                r.font.color.rgb = header_color; r.font.name = FONT
    for ri in range(1, len(tbl.rows)):
        for ci, cell in enumerate(tbl.rows[ri].cells):
            cell.fill.solid()
            cell.fill.fore_color.rgb = SOFT if (zebra and ri % 2 == 0) else WHITE
            cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
            cell.margin_left = Pt(6); cell.margin_right = Pt(6)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(body_size); r.font.name = FONT
                    r.font.color.rgb = INK
                    if ci == 0 and align_first_left:
                        r.font.bold = True; r.font.color.rgb = NAVY

def add_table(s, data, l, t, w, col_w=None, header_fill=NAVY, body_size=11.5,
              header_size=11.5, row_h=None):
    rows, cols = len(data), len(data[0])
    gtbl = s.shapes.add_table(rows, cols, l, t, w, Inches(0.4*rows)).table
    tblPr = gtbl._tbl.tblPr
    tblPr.set('firstRow', '0'); tblPr.set('bandRow', '0')
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = gtbl.cell(ri, ci)
            cell.text = ""
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = str(val); r.font.name = FONT
    if col_w:
        for ci, cw in enumerate(col_w):
            gtbl.columns[ci].width = cw
    if row_h:
        for ri in range(rows):
            gtbl.rows[ri].height = row_h
    style_table(gtbl, header_fill=header_fill, body_size=body_size, header_size=header_size)
    return gtbl

# ============================================================== SLIDE 1 — PORTADA
s = slide()
rect(s, 0, 0, SW, SH, fill=NAVY)
rect(s, 0, 0, SW, Inches(0.16), fill=GOLD)
tf = tb(s, Inches(0.9), Inches(1.55), Inches(11), Inches(0.4))
para(tf, "PROPUESTA DE ARRANQUE · CONFIDENCIAL", size=12, color=GOLD, bold=True, first=True)
tf = tb(s, Inches(0.9), Inches(2.1), Inches(11.4), Inches(2.0))
para(tf, "Tu ojo clínico,", size=40, color=WHITE, bold=True, first=True, leading=1.05)
para(tf, "vuelto un sistema que escala", size=40, color=WHITE, bold=True, leading=1.05)
tf = tb(s, Inches(0.9), Inches(4.15), Inches(10.8), Inches(1.4))
para(tf, "Ordenar y medir el motor de Reclutamiento & Selección y —con tu criterio— activar la IA y tu "
         "base de talento. Empezamos por una victoria concreta en 3 semanas, no por un gran compromiso.",
     size=16, color=RGBColor(0xC2,0xD4,0xE2), first=True, leading=1.3)
rect(s, Inches(0.9), Inches(6.25), Inches(11.5), Pt(1), fill=RGBColor(0x2A,0x44,0x5E))
for i,(a,b) in enumerate([("Preparado para","Virginia — CEO, Conexión Talento"),
                          ("Preparado por",f"{FIRM} · Data, IA & Transf. Digital"),
                          ("Fecha","Julio 2026")]):
    tf = tb(s, Inches(0.9+i*3.9), Inches(6.45), Inches(3.7), Inches(0.8))
    para(tf, a+":", size=10, color=RGBColor(0x9F,0xB4,0xC6), first=True, space_after=2)
    para(tf, b, size=11.5, color=WHITE, bold=True)

# ============================================================== SLIDE 2 — AGENDA
s = slide(); y = title_bar(s, "Hoja de ruta de hoy", "Agenda", 2)
items = [("01","El reto y el diagnóstico"),
         ("02","Nuestro enfoque: estandarizar antes de automatizar — con la IA, no sin ella"),
         ("03","Empecemos por una victoria: la Fase 0 y lo que te llevas en 3 semanas"),
         ("04","Cómo medimos el éxito y qué vale en tus dólares"),
         ("05","Inversión: tu decisión de hoy"),
         ("06","Riesgos, el camino completo y próximos pasos")]
for i,(n,t) in enumerate(items):
    yy = Inches(2.0 + i*0.78)
    rect(s, ML, yy, Inches(0.62), Inches(0.62), fill=SOFT, line=TEAL, line_w=Pt(1.2))
    tf = tb(s, ML, yy, Inches(0.62), Inches(0.62), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, n, size=15, color=TEAL, bold=True, align=PP_ALIGN.CENTER, first=True)
    tf = tb(s, Inches(1.5), yy, Inches(11.0), Inches(0.62), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, t, size=15.5, color=INK, first=True, leading=1.1)

# ============================================================== SLIDE 3 — EL RETO
s = slide(); y = title_bar(s, "01 · Entendimiento", "El reto", 3)
tf = tb(s, ML, y, CW, Inches(1.3))
para(tf, [("Conexión Talento vive de su motor de Reclutamiento & Selección — ", False, INK, 15),
          ("el 63% de los ingresos", True, NAVY, 15),
          (". Tiene activos poco comunes: el ojo clínico de Virginia, un ATS y ~4.000 candidatos. "
           "Lo que falta no son herramientas: es la ", False, INK, 15),
          ("columna vertebral", True, NAVY, 15),
          (" que vuelve ese talento escalable.", False, INK, 15)], first=True, leading=1.25)
box = rect(s, ML, Inches(3.05), CW, Inches(1.25), fill=RGBColor(0xEE,0xF5,0xF6))
rect(s, ML, Inches(3.05), Pt(4), Inches(1.25), fill=TEAL)
tf = tb(s, Inches(1.0), Inches(3.2), Inches(11.2), Inches(0.95), anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("La paradoja: ", True, NAVY, 15),
          ("vende “cercanía al candidato” mientras la critican en LinkedIn por falta de seguimiento, "
           "y ya presentó dos veces el mismo candidato al mismo cliente. El cuello de botella no es de "
           "manos — es de proceso.", False, INK, 15)], first=True, leading=1.25)
tf = tb(s, ML, Inches(4.7), CW, Inches(2.0))
for t in ["4 reclutadores trabajan de 4 formas distintas, sin documentar ni estandarizar.",
          "Cero métricas: sin línea base no hay ROI demostrable ni precio premium.",
          "La base de ~4.000 candidatos no es buscable — “el oro” está sin refinar."]:
    para(tf, t, size=14.5, color=INK, bullet=True, space_after=7, first=(t.startswith("4 ")), leading=1.15)

# ============================================================== SLIDE 4 — DIAGNÓSTICO
s = slide(); y = title_bar(s, "01 · Diagnóstico", "Dónde estás hoy", 4)
data = [["Dimensión","Estado actual","Implicación"],
        ["Proceso","4 reclutadores, 4 métodos; sin documentar","El criterio no escala ni se transfiere"],
        ["Datos","~4.000 candidatos sin etiquetar ni vigencia","No buscable; ~30–50% probablemente caduco"],
        ["Tecnología","Team Tailor con IA sin explotar","Se paga capacidad que no se usa"],
        ["Métricas","Inexistentes (“no métricas”)","Sin línea base no hay ROI ni precio premium"],
        ["Experiencia","Sin seguimiento; críticas en LinkedIn","La marca dice “cercanía”, el candidato no la siente"]]
add_table(s, data, ML, y, CW, col_w=[Inches(2.3), Inches(5.0), Inches(4.63)],
          body_size=12.5, header_size=12.5, row_h=Inches(0.72))

# ============================================================== SLIDE 5 — TESIS
s = slide(); rect(s,0,0,SW,SH,fill=NAVY); footer(s,5,light=True)
rect(s, 0, 0, Inches(0.16), SH, fill=GOLD)
tf = tb(s, ML, Inches(0.55), CW, Inches(0.4))
para(tf, "02 · NUESTRO ENFOQUE", size=11, color=GOLD, bold=True, first=True)
tf = tb(s, ML, Inches(1.35), Inches(11.8), Inches(2.4))
para(tf, [("Estandarizar ", True, WHITE, 30),("antes de automatizar.", True, TEALL, 30)],
     first=True, leading=1.1)
para(tf, "Meter IA sobre el desorden solo produce malos resultados más rápido. "
         "El diferenciador no es la IA — es la disciplina de proceso que ningún "
         "competidor copia sin el criterio humano detrás.",
     size=16, color=RGBColor(0xC2,0xD4,0xE2), space_before=10, leading=1.3)
steps = ["Estandarizar el criterio","Medir líneas base","Estructurar los datos","IA asistida"]
x = ML
for i,stp in enumerate(steps):
    w = Inches(2.7)
    rect(s, x, Inches(4.45), w, Inches(0.95), fill=RGBColor(0x16,0x3A,0x55), line=TEAL, line_w=Pt(1))
    tf = tb(s, x, Inches(4.45), w, Inches(0.95), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, f"{i+1}. {stp}", size=13.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER, first=True, leading=1.05)
    x = Emu(int(x) + int(w))
    if i < len(steps)-1:
        ar = tb(s, x, Inches(4.45), Inches(0.42), Inches(0.95), anchor=MSO_ANCHOR.MIDDLE)
        para(ar, "→", size=20, color=GOLD, bold=True, align=PP_ALIGN.CENTER, first=True)
        x = Emu(int(x) + int(Inches(0.42)))
tf = tb(s, ML, Inches(5.75), CW, Inches(1.0))
para(tf, [("Sí a la IA y a la velocidad que pediste — y para que rindan de verdad, primero el orden. ", True, GOLD, 15),
          ("La IA no es el final del camino; es la consecuencia de hacerlo bien.", False, RGBColor(0xC2,0xD4,0xE2), 15)],
     first=True, leading=1.3)

# ============================================================== SLIDE 6 — SÍ, Y
s = slide(); y = title_bar(s, "02 · Tu visión, afinada", "Sí, y — honramos lo que pediste", 6)
tf = tb(s, ML, Inches(1.62), CW, Inches(0.4))
para(tf, "Coincidimos en el qué. Aquí va el cómo, contigo como socia.", size=13.5, color=MUTED, first=True)
pairs = [
 ("“Quiero una persona más”","Sí, y primero una herramienta que multiplique las manos del equipo que ya tienes; si aun así hace falta gente, la sumas sabiendo exactamente para qué."),
 ("“Seis agentes de IA”","Sí, y arrancamos por el que más te quita el sueño y más valor te da. Uno bien hecho abre la puerta a los demás."),
 ("“La base es oro, vale ~60%”","Sí, y la puliremos para que valga lo que sabes: hoy es oro en bruto, sin clasificar y con una parte ya envejecida."),
 ("“Vendamos el benchmark salarial”","Sí, y por la puerta de adelante: anonimizado con dictamen legal, o de la mano de un proveedor licenciado."),
 ("“Salir del ATS / mismo día”","Sí, y primero exprimimos lo que ya pagas y lo medimos. Migrar después, con números. El “mismo día” es nuestro norte."),
]
yy = Inches(2.12)
for q, a in pairs:
    rect(s, ML, Emu(int(yy)+int(Inches(0.04))), Pt(3), Inches(0.5), fill=TEAL)
    tf = tb(s, Inches(0.92), yy, CW-Inches(0.22), Inches(0.6), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(q+"   ", True, NAVY, 13),("→  ", True, GOLD, 13),(a, False, INK, 13)], first=True, leading=1.12)
    yy = Emu(int(yy)+int(Inches(0.62)))
# blindar tu legado
by = Inches(5.55)
rect(s, ML, by, CW, Inches(1.0), fill=RGBColor(0xEE,0xF5,0xF6), line=TEAL, line_w=Pt(1))
rect(s, ML, by, Pt(4), Inches(1.0), fill=GOLD)
tf = tb(s, Inches(1.0), by, Inches(11.4), Inches(1.0), anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("Blindar tu legado:  ", True, NAVY, 14),
          ("no venimos a “extraerte el cerebro”, sino a poner tu ojo clínico de 20 años en un sistema que "
           "lo protege y lo multiplica — aunque tú no estés en la sala.", False, INK, 14)], first=True, leading=1.2)

# ============================================================== SLIDE 7 — EMPECEMOS POR UNA VICTORIA
s = slide(); y = title_bar(s, "03 · Tu decisión de hoy", "Empecemos por una victoria", 7)
# tarjeta héroe Fase 0
cardw = Inches(4.1)
rect(s, ML, y, cardw, Inches(4.55), fill=NAVY)
rect(s, ML, y, cardw, Inches(0.12), fill=GOLD)
tf = tb(s, Emu(int(ML)+int(Inches(0.3))), Emu(int(y)+int(Inches(0.35))), Inches(3.5), Inches(0.5))
para(tf, "FASE 0 · ARRANQUE", size=12, color=TEALL, bold=True, first=True)
tf = tb(s, Emu(int(ML)+int(Inches(0.3))), Emu(int(y)+int(Inches(0.9))), Inches(3.5), Inches(1.2))
para(tf, "US$2.500", size=44, color=WHITE, bold=True, first=True, leading=1.0)
tf = tb(s, Emu(int(ML)+int(Inches(0.3))), Emu(int(y)+int(Inches(2.0))), Inches(3.5), Inches(2.3))
for t in ["2–3 semanas","Factibilidad + 2 quick wins visibles","100% acreditable a la siguiente fase","Cerca de tu presupuesto"]:
    para(tf, t, size=13.5, color=RGBColor(0xC2,0xD4,0xE2), bullet=True, space_after=9,
         first=(t.startswith("2–3")), leading=1.15)
# bullets a la derecha
rx = Inches(5.15)
tf = tb(s, rx, y, Inches(7.4), Inches(3.7))
for i,t in enumerate(["Una prueba pagada, con resultado visible en 2–3 semanas.",
          "Cerca de tu presupuesto, no un salto al vacío.",
          "El 100% se acredita a la siguiente fase si avanzas — baja tu riesgo, no es descuento.",
          "Estandarizar, IA y escala los cotizamos DESPUÉS, con tus datos reales en mano.",
          "No te pedimos comprometerte con todo a ciegas: primero ves, luego decides."]):
    para(tf, t, size=15, color=INK, bullet=True, space_after=12, first=(i==0), leading=1.2)
# cierre
cy = Inches(6.05)
rect(s, rx, cy, Inches(7.4), Inches(0.7), fill=SOFT, line=GOLD, line_w=Pt(1))
tf = tb(s, Emu(int(rx)+int(Inches(0.2))), cy, Inches(7.0), Inches(0.7), anchor=MSO_ANCHOR.MIDDLE)
para(tf, "Tu decisión de hoy es pequeña; el resto se gana con datos.", size=14, color=NAVY, bold=True, first=True)

# ============================================================== SLIDE 8 — LO QUE TE LLEVAS (DEMO VIVA)
s = slide(); y = title_bar(s, "03 · La demo viva", "Lo que te llevas en 3 semanas", 8)
tf = tb(s, ML, y, CW, Inches(0.45))
para(tf, [("Con TUS datos, no en abstracto: lo ves, lo tocas, y desde ahí decides el siguiente paso.", True, TEAL, 14.5)],
     first=True, leading=1.15)
items = [
 ("La IA trabajando con tus datos","20–30 CVs reales de una vacante tuya abierta, ordenados con tu criterio replicado — en un entorno con DPA, cumpliendo desde el día 1."),
 ("Tu CV con el sello de la firma","plantilla única, ya montada sobre un candidato real."),
 ("Registro anti-duplicado funcionando","se acaba el bochorno de presentar dos veces al mismo candidato."),
 ("El semáforo honesto","te decimos si tu base realmente se puede exportar y cuánto de los 4.000 sigue vigente."),
]
yy = Inches(2.45)
for h,t in items:
    rect(s, ML, yy, Inches(0.5), Inches(0.5), fill=TEAL)
    tf = tb(s, ML, yy, Inches(0.5), Inches(0.5), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, "✓", size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
    tf = tb(s, Inches(1.45), yy, Inches(11.0), Inches(0.95), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(h+":  ", True, NAVY, 15),(t, False, INK, 15)], first=True, leading=1.18)
    yy = Emu(int(yy)+int(Inches(0.98)))

# ============================================================== SLIDE 9 — MÉTRICAS
s = slide(); y = title_bar(s, "04 · Cómo medimos el éxito", "De “cero métricas” a un cuadro de mando", 9)
data = [["KPI","Qué mide","Línea base","Meta"],
        ["Time-to-terna","Días de apertura a terna entregada","~5 días","48–72h"],
        ["Aceptación de terna","% de ternas que el cliente aprueba","a medir","↑ sostenido"],
        ["Rotación a 90 días","% de colocados que salen <90 días","a medir","↓"],
        ["Cobertura base interna","% de mandatos cubiertos sin sourcing externo","~0% hoy","↑"],
        ["Salud de la base","% de candidatos vigentes y etiquetados","desconocida","↑"],
        ["Experiencia (CSAT)","% de candidatos satisfechos (con n)","a medir","≥ 80%"]]
add_table(s, data, ML, y, CW, col_w=[Inches(2.6),Inches(5.7),Inches(2.0),Inches(1.63)],
          body_size=11.5, header_size=11.5, row_h=Inches(0.6))
tf = tb(s, ML, Inches(6.5), CW, Inches(0.4))
para(tf, "Crear las líneas base (hoy inexistentes) es parte del valor de la Fase 1 — y la prueba de tu ROI.",
     size=11, color=MUTED, first=True)

# ============================================================== SLIDE 10 — ROI
s = slide(); y = title_bar(s, "04 · El retorno en tus números", "Lo que vale cada terna entregada más rápido", 10)
tf = tb(s, ML, y, CW, Inches(0.7))
para(tf, "Pasar tu terna de ~5 días a 48–72h no es solo comodidad: es ganar mandatos que hoy se enfrían "
         "o se lleva otro.", size=14.5, color=INK, first=True, leading=1.2)
# fórmula
fy = Inches(2.55)
rect(s, ML, fy, CW, Inches(0.92), fill=NAVY)
rect(s, ML, fy, Pt(4), Inches(0.92), fill=GOLD)
tf = tb(s, Inches(1.0), fy, Inches(11.4), Inches(0.92), anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("Fórmula:  ", True, TEALL, 15),
          ("(colocaciones extra que cierras por presentar primero) × (tu fee por colocación) = ", False, WHITE, 15),
          ("retorno", True, GOLD, 15),
          (".  Una sola colocación adicional ya cubre la Fase 0.", False, WHITE, 15)], first=True, leading=1.2)
# ejemplo
ey = Inches(3.75)
rect(s, ML, ey, CW, Inches(1.0), fill=RGBColor(0xEE,0xF5,0xF6), line=TEAL, line_w=Pt(1))
tf = tb(s, Inches(1.0), ey, Inches(11.4), Inches(1.0), anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("Ejemplo — lo llenamos con tus números:  ", True, NAVY, 13.5),
          ("salario anual US$18,000 × fee 15% = US$2,700 por colocación. Si esto te ayuda a cerrar 1 "
           "colocación que hoy perderías, ya recuperaste los US$2,500 de la Fase 0.  ", False, INK, 13.5),
          ("(Números de marcador, no promesa.)", False, MUTED, 12.5)], first=True, leading=1.2)
# honesto + caveat
tf = tb(s, ML, Inches(5.0), CW, Inches(1.5))
para(tf, [("Lo honesto: ", True, NAVY, 13.5),
          ("todavía no tengo tu fee promedio ni cuántos mandatos manejas al mes. Mándamelos HOY por "
           "WhatsApp y traigo tu ROI con tu número real al viernes —no un ejemplo—.", False, INK, 13.5)],
     bullet=True, space_after=8, first=True, leading=1.2)
para(tf, [("Caveat: ", True, NAVY, 13.5),
          ("la velocidad te pone primera en la fila; el cierre sigue dependiendo de tu ojo y del cliente. "
           "Buscamos ponerte primera en la fila — no prometemos colocaciones.", False, INK, 13.5)],
     bullet=True, space_after=8, leading=1.2)

# ============================================================== SLIDE 11 — INVERSIÓN
s = slide(); y = title_bar(s, "05 · Inversión", "Tu decisión de hoy — y el horizonte", 11)
# héroe F0
rect(s, ML, y, CW, Inches(1.4), fill=RGBColor(0xEE,0xF5,0xF6), line=TEAL, line_w=Pt(1.5))
rect(s, ML, y, Inches(0.14), Inches(1.4), fill=TEAL)
tf = tb(s, Inches(1.05), Emu(int(y)+int(Inches(0.18))), Inches(7.8), Inches(1.1), anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("Empezamos aquí — Fase 0", True, NAVY, 17)], first=True, space_after=3)
para(tf, "Factibilidad pagada + 2 quick wins · 2–3 semanas · 100% acreditable a la siguiente fase",
     size=12.5, color=MUTED, space_after=3)
para(tf, "Tu seguro: los activos tangibles son tuyos aunque la Puerta 0 salga roja y no continúes.",
     size=12, color=TEAL, bold=True)
tf = tb(s, Inches(9.2), y, Inches(3.3), Inches(1.4), anchor=MSO_ANCHOR.MIDDLE)
para(tf, "US$2.500", size=34, color=TEAL, bold=True, align=PP_ALIGN.RIGHT, first=True)
# ancla de la cuenta real (temp + F0 vs 5o reclutador)
tf = tb(s, ML, Inches(3.16), CW, Inches(0.3))
para(tf, [("La cuenta real:  ", True, GOLD, 12),
          ("temporal (~$2.000) + Fase 0 ($2.500) = $4.500  ·  vs. el 5º reclutador que pediste (~$15k/año). Esa es la vara.", False, MUTED, 12)],
     first=True, leading=1.1)
# horizonte (de-enfatizado)
tf = tb(s, ML, Inches(3.56), CW, Inches(0.35))
para(tf, "El horizonte completo (referencia — no es lo que decides hoy; cada fase se cotiza y firma por separado):",
     size=12.5, color=MUTED, first=True)
hz = [("Sistema","Fases 0 + 1","~US$12.000"),
      ("Sistema + piloto IA","Fases 0 + 1 + 2","~US$24.000"),
      ("Programa completo","Fases 0–3","desde ~US$39.000")]
cw = Inches(3.83)
for i,(t,sub,price) in enumerate(hz):
    x = Emu(int(ML) + i*int(Inches(4.05)))
    rect(s, x, Inches(3.95), cw, Inches(1.35), fill=SOFT, line=LINE, line_w=Pt(1))
    rect(s, x, Inches(3.95), cw, Inches(0.1), fill=NAVY)
    tf = tb(s, Emu(int(x)+int(Inches(0.2))), Inches(4.15), Inches(3.4), Inches(1.1))
    para(tf, t, size=13, color=NAVY, bold=True, first=True, space_after=2)
    para(tf, sub, size=11, color=MUTED, space_after=6)
    para(tf, price, size=17, color=MUTED, bold=True)
# modelo de cobro
my = Inches(5.65)
rect(s, ML, my, CW, Inches(0.85), fill=NAVY)
tf = tb(s, Inches(1.0), my, Inches(11.4), Inches(0.85), anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("Cómo cobramos:  ", True, GOLD, 13.5),
          ("precio cerrado por fase — pagas solo lo que apruebas en cada puerta. Sin retainer ni cobro por "
           "“éxito” al inicio: no hay aún línea base para medirlo de forma justa.", False, WHITE, 13.5)],
     first=True, leading=1.2)

# ============================================================== SLIDE 12 — RIESGOS
s = slide(); y = title_bar(s, "06 · Riesgos y cumplimiento", "Lo gestionamos por diseño", 12)
data = [["Riesgo","Mitigación"],
        ["Protección de datos — El Salvador ya regula (Decreto 144/2024, vigente)","Consentimiento por finalidad, registro, retención y responsable de datos desde la Fase 1. Cumplimiento por diseño = argumento de venta"],
        ["Vender el benchmark salarial = bandera roja","Congelado hasta dictamen de abogado salvadoreño + anonimización irreversible (k-anonimato)"],
        ["Sesgo en cribado automático","Humano-en-el-bucle obligatorio; la IA nunca rechaza sola; auditoría de sesgo y trazabilidad"],
        ["API de Team Tailor no verificada","Spike de factibilidad en la Fase 0 (semana 1); plan B de captura acotada"],
        ["Datos en Guatemala (transferencia SV–GT)","R&S opera en SV+GT. GT aún sin ley general; transferencia con DPA y minimización. El estándar se construye una vez y se despliega en ambos países"],
        ["Adopción del equipo","Co-diseño con los reclutadores; el rol extra como embajador de marca, no operativo"]]
add_table(s, data, ML, y, CW, col_w=[Inches(4.6),Inches(7.33)],
          body_size=10.5, header_size=12, row_h=Inches(0.66))

# ============================================================== SLIDE 13 — HORIZONTE (FASES OVERVIEW)
s = slide(); y = title_bar(s, "06 · El camino completo", "Las cuatro fases — el mapa, no el pedido de hoy", 13)
# banner de referencia
rect(s, ML, y, CW, Inches(0.55), fill=RGBColor(0xEE,0xF5,0xF6), line=GOLD, line_w=Pt(1))
tf = tb(s, Inches(1.0), y, Inches(11.4), Inches(0.55), anchor=MSO_ANCHOR.MIDDLE)
para(tf, [("Referencia: ", True, NAVY, 12.5),
          ("hoy solo decides la Fase 0. Cada fase siguiente se cotiza y se firma por separado, con datos reales en mano.", False, INK, 12.5)],
     first=True, leading=1.1)
yy = Emu(int(y)+int(Inches(0.72)))
data = [["Fase","Nombre","Duración","Esfuerzo","Inversión","Resultado / Puerta"],
        ["0 ◀","Factibilidad + Quick Wins","2–3 sem","8–12 jorn.","US$2.500*","Empezamos aquí: semáforo + 2 quick wins"],
        ["1","Estandarizar, Documentar, Medir","4–6 sem","20–30 jorn.","~US$9.500","La “columna vertebral”: SOPs, rúbrica, métricas"],
        ["2","IA asistida (modular)","6–8 sem","25–35 jorn.","~US$12.000","Screener/Ranker + Generador de CV (con HITL)"],
        ["3","Activo de datos y escala","horizonte","a dimensionar","desde 15.000 /\nretainer 3.500/mes","Base buscable, mercado interno, benchmark"]]
add_table(s, data, ML, yy, CW,
          col_w=[Inches(0.85),Inches(2.95),Inches(1.2),Inches(1.3),Inches(1.85),Inches(3.78)],
          body_size=11, header_size=11, row_h=Inches(0.78))
tf = tb(s, ML, Inches(6.55), CW, Inches(0.4))
para(tf, "* La Fase 0 acredita el 100% a la Fase 1. Cifras indicativas en USD, a calibrar con tu tarifa local y tu fee por colocación.",
     size=9.5, color=MUTED, first=True)

# ============================================================== SLIDE 14 — ROADMAP
s = slide(); y = title_bar(s, "06 · Roadmap", "Cómo se ve en el tiempo", 14)
chart_x = Inches(2.35); chart_w = Inches(10.0); weeks = 16
wk = Emu(int(chart_w) // weeks)
for i in range(0, weeks+1, 2):
    gx = Emu(int(chart_x) + i*int(wk))
    rect(s, gx, Inches(2.05), Pt(0.8), Inches(3.4), fill=LINE)
    lab = tb(s, Emu(gx-int(Inches(0.2))), Inches(1.78), Inches(0.4), Inches(0.25))
    para(lab, f"S{i if i>0 else 1}", size=8.5, color=MUTED, align=PP_ALIGN.CENTER, first=True)
bars = [("Fase 0", 0, 3, TEALL),("Fase 1", 2, 8, TEAL),
        ("Fase 2", 6, 13, DEEP),("Fase 3", 10, 16, GOLD)]
for i,(name, a, b, col) in enumerate(bars):
    ry = Inches(2.25 + i*0.78)
    lab = tb(s, ML, ry, Inches(1.6), Inches(0.5), anchor=MSO_ANCHOR.MIDDLE)
    para(lab, name, size=13, color=NAVY, bold=True, first=True)
    bx = Emu(int(chart_x) + a*int(wk)); bw = Emu((b-a)*int(wk))
    rect(s, bx, ry, bw, Inches(0.5), fill=col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
tf = tb(s, ML, Inches(5.7), CW, Inches(1.1))
for t,fr in [("Hoy decides la Fase 0 (semanas 1–3). Lo demás se cotiza al pasar cada puerta.", True),
             ("Las fases se solapan: la extracción de datos corre en paralelo a la documentación.", False),
             ("Cada puerta go/no-go es una decisión explícita de continuar, reordenar o detener.", False)]:
    para(tf, t, size=13, color=INK, bullet=True, space_after=6, first=fr, leading=1.15)

# ============================================================== SLIDES 15-17 — ANEXO · DETALLE FASES
def phase_slide(n, tag, meta, objetivo, entregables, puerta, color, idx, kicker=None):
    s = slide(); y = title_bar(s, kicker or f"Anexo · Fase {n}", tag, idx)
    rect(s, ML, y, Inches(2.0), Inches(0.42), fill=color)
    tf = tb(s, ML, y, Inches(2.0), Inches(0.42), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, f"  FASE {n}", size=13, color=WHITE, bold=True, first=True)
    tf = tb(s, Inches(2.9), y, Inches(9.7), Inches(0.42), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, meta, size=11.5, color=MUTED, align=PP_ALIGN.RIGHT, first=True)
    yy = Emu(int(y)+int(Inches(0.62)))
    tf = tb(s, ML, yy, CW, Inches(0.6))
    para(tf, [("Objetivo:  ", True, NAVY, 14),(objetivo, False, INK, 14)], first=True, leading=1.15)
    tf = tb(s, ML, Emu(int(yy)+int(Inches(0.62))), CW, Inches(3.0))
    para(tf, "Entregables", size=13, color=TEAL, bold=True, first=True, space_after=5)
    for i,e in enumerate(entregables):
        if isinstance(e, tuple):
            para(tf, [(e[0]+": ", True, NAVY, 12.5),(e[1], False, INK, 12.5)], bullet=True, space_after=5, leading=1.12)
        else:
            para(tf, e, size=12.5, color=INK, bullet=True, space_after=5, leading=1.12)
    by = Inches(6.05)
    rect(s, ML, by, CW, Inches(0.75), fill=SOFT, line=color, line_w=Pt(1))
    rect(s, ML, by, Pt(4), Inches(0.75), fill=color)
    tf = tb(s, Inches(0.95), by, Inches(11.5), Inches(0.75), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, [(f"Puerta {n} →  ", True, color, 12),(puerta, False, INK, 12)], first=True, leading=1.12)

phase_slide("1","Estandarizar, Documentar y Medir","4–6 semanas · 20–30 jornadas · ~US$9.500",
    "instalar la columna vertebral. Aquí vive el valor real y la capacidad que se queda.",
    [("Mapa de proceso “deber ser”","los 8 pasos con RACI, SLA por etapa y dueño nombrado"),
     ("SOP de cribado + rúbrica + golden set","tu ojo clínico hecho criterios y pesos — el entregable estrella"),
     ("Scorecard de intake","1 página por vacante: must/nice-to-have, pesos, éxito a 90 días"),
     ("Árbol de métricas + estandarización del ATS","1 North Star + ~10 KPIs; etapas y campos obligatorios"),
     ("Gobernanza de datos + taxonomía v1","RoPA-lite, base legal por finalidad, diccionario de skills")],
    "etapas estandarizadas; SOP + rúbrica firmados; golden set armado; baseline medido. La IA no arranca sin esto.",
    TEAL, 15, kicker="Anexo · Horizonte — Fase 1")

phase_slide("2","IA asistida (modular)","6–8 semanas · 25–35 jornadas · ~US$12.000",
    "acelerar lo que ya probó su criterio. Solo módulos que pasen la “Puerta Listo para IA”.",
    [("Screener/Ranker asistido (~US$7.000)","des-identifica PII → puntúa vs scorecard con justificación → pre-ordena. El humano valida; la IA nunca rechaza sola"),
     ("Generador de CV con branding (~US$4.000)","extracción a esquema fijo + plantilla determinista"),
     ("Medición vs línea base + auditoría de sesgo","trazabilidad de prompts, criterios y scores")],
    "concordancia IA-vs-ojo clínico ≥ umbral en búsquedas reales; sin sesgo relevante; mejora medible de time-to-terna.",
    DEEP, 16, kicker="Anexo · Horizonte — Fase 2")

phase_slide("3","Activo de datos y escala","horizonte · land-and-expand · desde US$15.000 o retainer US$3.500/mes",
    "convertir la base en activo líquido y abrir nuevas líneas. Condicionada a las fases previas.",
    [("Enriquecimiento de la base","parsing, golden record, score de vigencia → búsqueda por skill e internal-fill"),
     ("Mercado interno / historial de candidato","+ portal de auto-actualización"),
     ("Decisión augmentar-vs-migrar Team Tailor","con datos y caso de negocio"),
     ("Benchmark de compensaciones","solo con validación legal, k-anonimato y consentimiento — o alianza con proveedor licenciado")],
    "cada nueva línea se autofinancia con el valor demostrado en las fases previas.",
    GOLD, 17, kicker="Anexo · Horizonte — Fase 3")

# ============================================================== SLIDE 18 — NECESITAMOS / PASOS
s = slide(); y = title_bar(s, "06 · Para arrancar", "Qué necesitamos de ti y próximos pasos", 18)
tf = tb(s, ML, y, Inches(6.0), Inches(0.4))
para(tf, "Qué necesitamos de Conexión Talento", size=14, color=TEAL, bold=True, first=True)
tf = tb(s, ML, Emu(int(y)+int(Inches(0.5))), Inches(6.0), Inches(3.8))
for t in ["Un CV real anonimizado (por correo) — para el teaser del viernes.",
          "Acceso de lectura a Team Tailor — el lunes (lead time).",
          "Contacto de tu abogado/a — el lunes.",
          "Fee por colocación + volumen de mandatos/mes.",
          "Confirmación del consentimiento de la base."]:
    para(tf, t, size=13, color=INK, bullet=True, space_after=8, first=(t.startswith("Un CV")), leading=1.15)
rx = Inches(7.1)
tf = tb(s, rx, y, Inches(5.5), Inches(0.4))
para(tf, "Próximos pasos", size=14, color=TEAL, bold=True, first=True)
for i,t in enumerate(["Aprobar la Fase 0 (US$2.500) — tu decisión de hoy.",
                      "Hoy: CV anonimizado + tu fee por WhatsApp → teaser branded + registro anti-duplicado en vivo esta semana.",
                      "Lunes: acceso a Team Tailor + abogado/a. Demo completa + semáforo en 2–3 semanas."]):
    yy = Emu(int(y)+int(Inches(0.55)) + i*int(Inches(1.0)))
    rect(s, rx, yy, Inches(0.55), Inches(0.55), fill=TEAL)
    tf = tb(s, rx, yy, Inches(0.55), Inches(0.55), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, str(i+1), size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER, first=True)
    tf = tb(s, Emu(int(rx)+int(Inches(0.8))), yy, Inches(4.6), Inches(0.85), anchor=MSO_ANCHOR.MIDDLE)
    para(tf, t, size=13, color=INK, first=True, leading=1.15)

# ============================================================== SLIDE 19 — CIERRE
s = slide(); rect(s,0,0,SW,SH,fill=NAVY)
rect(s, 0, 0, SW, Inches(0.16), fill=GOLD)
tf = tb(s, Inches(0.9), Inches(2.6), Inches(11.6), Inches(2.0))
para(tf, "Empezamos pequeño, demostramos valor con tus datos,", size=27, color=WHITE, bold=True, first=True, leading=1.14)
para(tf, "y escalamos solo lo que funciona.", size=27, color=TEALL, bold=True, leading=1.14)
tf = tb(s, Inches(0.9), Inches(4.45), Inches(11), Inches(1.0))
para(tf, "Tu decisión de hoy es chica: una victoria visible en 3 semanas. El resto se gana con datos, "
         "fase por fase — y con tu ojo clínico en el centro de todo.",
     size=15, color=RGBColor(0xC2,0xD4,0xE2), first=True, leading=1.3)
tf = tb(s, Inches(0.9), Inches(6.4), Inches(11.6), Inches(0.5))
para(tf, f"Gracias, Virginia. — {FIRM} · Data, IA & Transformación Digital",
     size=12, color=GOLD, bold=True, first=True)

prs.save("Presentacion-Conexion-Talento.pptx")
print("PPTX guardado:", len(prs.slides._sldIdLst), "slides")
